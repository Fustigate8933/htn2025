from pptx import Presentation
from PyPDF2 import PdfReader
from services.generate_speech import generate_speech
from google.cloud import storage
import io
import time


class PPTProcessor:
    def __init__(self):
        # Initialize the official GCS client
        self.gcs = storage.Client()

    def extract_ppt_text(self, ppt_path):
        bucket_name = "hack-the-north-bucket"

        # Fetch file from GCS if it's a cloud path, otherwise read local file
        if ppt_path.startswith("gs://"):
            # format: gs://bucket/path/to/file.pptx
            _, _, bucket_name, *blob_parts = ppt_path.split("/")
            blob_name = "/".join(blob_parts)
            bucket = self.gcs.bucket(bucket_name)
            blob = bucket.blob(blob_name)
            data = blob.download_as_bytes()
        else:
            # local file
            with open(ppt_path, "rb") as f:
                data = f.read()

        prs = Presentation(io.BytesIO(data))
        slides_text = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            if slide_content:
                slides_text.append({
                    "page_num": i,
                    "text": "\n".join(slide_content)
                })
        return slides_text

    def extract_pdf_text(self, pdf_path):
        text = []
        with open(pdf_path, "rb") as f:
            reader = PdfReader(f)
            for i, page in enumerate(reader.pages, start=1):
                content = page.extract_text()
                if content and content.strip():
                    text.append({
                        "page_num": i,
                        "text": content.strip()
                    })
        return text

    def build_prompt(self, ppt_text):
        """
        Build the prompt to send to cohere API based on the extracted text
        """
        prompt = (
            f"Use English to generate a speech draft for Slide/Page {ppt_text['page_num']}.\n\n"
            f"Content:\n{ppt_text['text']}\n\n"
            "Please output a clear and complete speech text for this slide/page. "
            "Do not summarize, make it ready to speak."
        )
        return prompt

    def file_to_speech(self, file_path: str) -> dict:
        if file_path.endswith(".pptx"):
            pages = self.extract_ppt_text(file_path)
        elif file_path.endswith(".pdf"):
            pages = self.extract_pdf_text(file_path)
        else:
            raise ValueError("Unsupported file format: Only .pptx and .pdf are supported")

        result = {}
        for page in pages:
            print(f"Processing Page {page['page_num']}...")
            prompt = self.build_prompt(page)
            speech_text = generate_speech(prompt)
            result[page["page_num"]] = speech_text
            time.sleep(7)  # avoid API rate limits

        return result


if __name__ == "__main__":
    ppt_processor = PPTProcessor()

    ppt_path = "/mnt/ianch-Secondary/Downloads/tmp/CS136W25 Midterm Review .pptx"

    ppt_text = ppt_processor.extract_ppt_text(ppt_path)

    for page in ppt_text:
        prompt = ppt_processor.build_prompt(page)
        speech_text = generate_speech(prompt)
        print(f"Generated speech for page {page['page_num']}:\n{speech_text}\n")

